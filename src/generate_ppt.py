from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_stylized_presentation():
    prs = Presentation()
    
    # --- DESIGN CONSTANTS ---
    # Color Palette: Dark Cinematic Theme
    BG_COLOR = RGBColor(20, 20, 20)      # Almost Black
    TITLE_COLOR = RGBColor(229, 9, 20)   # Netflix Red
    TEXT_COLOR = RGBColor(245, 245, 245) # Off-White
    ACCENT_COLOR = RGBColor(80, 80, 80)  # Dark Grey accent

    def set_slide_background(slide):
        """Sets the background color of a slide to BG_COLOR."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_header_bar(slide):
        """Adds a decorative red line at the top."""
        left = Inches(0)
        top = Inches(1.3)
        width = prs.slide_width
        height = Inches(0.05)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = TITLE_COLOR
        shape.line.fill.background() # No outline

    def slide_title_style(shape, text):
        shape.text = text
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.font.name = 'Arial Black'
        paragraph.font.size = Pt(40)
        paragraph.font.color.rgb = TITLE_COLOR
        paragraph.alignment = PP_ALIGN.LEFT

    def slide_body_style(shape, content_list):
        tf = shape.text_frame
        tf.clear() # Clear default empty paragraph
        
        for item in content_list:
            p = tf.add_paragraph()
            p.text = item
            p.font.name = 'Calibri'
            p.font.size = Pt(24)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(14)
            p.space_before = Pt(6)

    # --- SLIDES GENERATION ---

    def add_content_slide(title_text, content_list):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide)
        add_header_bar(slide)
        
        # Title
        slide_title_style(slide.shapes.title, title_text)
        
        # Body
        body_shape = slide.shapes.placeholders[1]
        slide_body_style(body_shape, content_list)

    # 1. Main Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    
    title.text = "CINE CAPSULE"
    title.text_frame.paragraphs[0].font.name = 'Arial Black'
    title.text_frame.paragraphs[0].font.size = Pt(60)
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    subtitle.text = "Sentiment-Driven Recommendations & Gamification\n\nDepartment of Computer Science"
    subtitle.text_frame.paragraphs[0].font.name = 'Arial'
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR

    # Add a decorative rectangle behind title for effect
    left = Inches(1)
    top = Inches(2.8)
    width = Inches(8)
    height = Inches(0.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.fill.background()

    # 2. The Problem
    add_content_slide("The Problem", [
        "⚠️ Information Overload: Viewers drown in choices.",
        "🚫 Flawed Ratings: 7/10 tells you 'how much', not 'why'.",
        "📉 Low Retention: Users leave after watching 1 movie.",
        "Need: A system that finds quality (AI) and rewards loyalty (Gamification)."
    ])

    # 3. The Solution
    add_content_slide("The Solution", [
        "A Dual-Module Architecture:",
        "1. Sentiment AI: Trusts reviews over ratings.",
        "2. Gamification: Rewards 'Watch Behavior'.",
        "",
        "\"Don't just watch movies. Experience them.\""
    ])

    # 4. System Logic
    add_content_slide("System Architecture", [
        "INPUT: 50,000 IMDB Reviews (Raw Text).",
        "↓",
        "PROCESSING: NLP Cleaning & TF-IDF Vectorization.",
        "↓",
        "AI ENGINE: Logistic Regression (Positive/Negative).",
        "↓",
        "OUTPUT: Top 5 Hybrid-Ranked Movies."
    ])

    # 5. Recommendation Core
    add_content_slide("Module 1: AI Engine", [
        "The Hybrid Formula:",
        "Score = (Sentiment × 0.7) + (Rating × 0.3)",
        "",
        "Why 70% Sentiment?",
        "• A 5-star rating is subjective.",
        "• A generic 'Amazing movie' review is objective evidence of joy."
    ])

    # 6. Gamification Core
    add_content_slide("Module 2: Gamification", [
        "The 'CineCapsule' Engagement Loop:",
        "• Watch Movie → System Logs Genre/Time.",
        "• Check conditions against 40+ Badges.",
        "• Unlock Achievement → Dopamine Hit.",
        "",
        "Examples:",
        "🏆 'Night Owl' (Watch after 11 PM)",
        "🏆 'Action Rookie' (Watch 3 Action films)"
    ])

    # 7. Tech Stack
    add_content_slide("Technology Stack", [
        "🐍 Python 3.8: Core Logic",
        "🐼 Pandas & NumPy: Data Processing",
        "🤖 Scikit-Learn: Machine Learning Models",
        "🗣️ NLTK: Natural Language Processing",
        "📊 Matplotlib: Data Visualization"
    ])

    # 8. Demo Output
    add_content_slide("Live Demo Results", [
        "✔ Sentiment Accuracy: 89.4%",
        "✔ Recommendation: Successfully filters low-quality hits.",
        "✔ Gamification: Real-time badge unlocking works flawlessly.",
        "",
        "[Live Demo Executed via src/run_achievement_demo.py]"
    ])

    # 9. Conclusion
    add_content_slide("Conclusion", [
        "CineCapsule is not just a Recommender.",
        "It is an Ecosystem.",
        "",
        "By merging Sentiment Intelligence with Gamification strategies, we solve Discovery AND Retention simultaneously."
    ])

    # Save
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    output_path = os.path.join(base_dir, 'reports', 'CineCapsule_Premium_Presentation.pptx')
    prs.save(output_path)
    print(f"Stylized presentation saved to: {output_path}")

if __name__ == "__main__":
    create_stylized_presentation()
